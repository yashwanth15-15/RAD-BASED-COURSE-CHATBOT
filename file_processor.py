from langchain_community.document_loaders import PyPDFLoader, TextLoader
from langchain_text_splitters import CharacterTextSplitter
from langchain_core.documents import Document
from pptx import Presentation
import tempfile


def process_files(files):
    documents = []

    for file in files:
        with tempfile.NamedTemporaryFile(delete=False) as tmp:
            tmp.write(file.read())
            path = tmp.name

        # ---------- PDF ----------
        if file.name.endswith(".pdf"):
            loader = PyPDFLoader(path)
            docs = loader.load()

            for d in docs:
                documents.append(
                    Document(
                        page_content=d.page_content,
                        metadata={"source": file.name}
                    )
                )

        # ---------- TXT ----------
        elif file.name.endswith(".txt"):
            loader = TextLoader(path)
            docs = loader.load()

            for d in docs:
                documents.append(
                    Document(
                        page_content=d.page_content,
                        metadata={"source": file.name}
                    )
                )

        # ---------- PPT ----------
        elif file.name.endswith(".pptx"):
            prs = Presentation(path)

            for i, slide in enumerate(prs.slides):
                text = ""

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + " "

                if text.strip():
                    documents.append(
                        Document(
                            page_content=text,
                            metadata={
                                "source": file.name,
                                "type": "ppt",
                                "slide": i + 1
                            }
                        )
                    )

    splitter = CharacterTextSplitter(
        chunk_size=600,
        chunk_overlap=100
    )

    return splitter.split_documents(documents)
